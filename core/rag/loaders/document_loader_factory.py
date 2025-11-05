"""
Document Loader Factory
Factory 패턴: 파일 형식별 로더 생성
"""

from typing import List, Optional
from pathlib import Path
from langchain.schema import Document
from core.logging import get_logger

logger = get_logger("document_loader_factory")


class DocumentLoaderFactory:
    """문서 로더 팩토리"""
    
    @staticmethod
    def load_document(file_path: str) -> List[Document]:
        """
        Load document based on file type
        
        Args:
            file_path: File path
            
        Returns:
            List of Document objects
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_type = path.suffix.lower()
        
        # 파일 형식별 로더 선택
        if file_type == '.pdf':
            return DocumentLoaderFactory._load_pdf(path)
        elif file_type == '.docx':
            return DocumentLoaderFactory._load_word(path)
        elif file_type in ['.xlsx', '.xls']:
            return DocumentLoaderFactory._load_excel(path)
        elif file_type == '.pptx':
            return DocumentLoaderFactory._load_powerpoint(path)
        elif file_type in ['.hwp', '.hwpx']:
            logger.warning(f"HWP files not supported: {file_type}")
            return []
        elif file_type == '.csv':
            return DocumentLoaderFactory._load_csv(path)
        elif file_type == '.txt':
            return DocumentLoaderFactory._load_text(path)
        elif file_type == '.json':
            return DocumentLoaderFactory._load_json(path)
        elif file_type in ['.png', '.jpg', '.jpeg']:
            return DocumentLoaderFactory._load_image(path)
        else:
            logger.warning(f"Unsupported file type: {file_type}")
            return []
    
    @staticmethod
    def _load_pdf(path: Path) -> List[Document]:
        """Load PDF file"""
        try:
            from PyPDF2 import PdfReader
            
            reader = PdfReader(str(path))
            documents = []
            
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    doc = Document(
                        page_content=text,
                        metadata={"source": str(path)}
                    )
                    documents.append(doc)
            
            logger.info(f"Loaded PDF: {path.name} ({len(documents)} pages)")
            return documents
            
        except Exception as e:
            logger.error(f"Failed to load PDF: {e}")
            return []
    
    @staticmethod
    def _load_word(path: Path) -> List[Document]:
        """Load Word file"""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(str(path))
            text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            
            if text.strip():
                document = Document(
                    page_content=text,
                    metadata={"source": str(path)}
                )
                logger.info(f"Loaded Word: {path.name}")
                return [document]
            
            return []
            
        except Exception as e:
            logger.error(f"Failed to load Word: {e}")
            return []
    
    @staticmethod
    def _load_excel(path: Path) -> List[Document]:
        """Load Excel file"""
        try:
            import pandas as pd
            
            # 모든 시트 읽기
            excel_file = pd.ExcelFile(str(path))
            documents = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # DataFrame을 텍스트로 변환
                text = f"Sheet: {sheet_name}\n\n"
                text += df.to_string(index=False)
                
                doc = Document(
                    page_content=text,
                    metadata={"source": str(path)}
                )
                documents.append(doc)
            
            logger.info(f"Loaded Excel: {path.name} ({len(documents)} sheets)")
            return documents
            
        except Exception as e:
            logger.error(f"Failed to load Excel: {e}")
            return []
    
    @staticmethod
    def _load_csv(path: Path) -> List[Document]:
        """Load CSV file"""
        try:
            import pandas as pd
            
            df = pd.read_csv(str(path))
            text = df.to_string(index=False)
            
            doc = Document(
                page_content=text,
                metadata={"source": str(path)}
            )
            
            logger.info(f"Loaded CSV: {path.name}")
            return [doc]
            
        except Exception as e:
            logger.error(f"Failed to load CSV: {e}")
            return []
    
    @staticmethod
    def _load_text(path: Path) -> List[Document]:
        """Load text file"""
        encodings = ['utf-8', 'euc-kr', 'cp949', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(path, 'r', encoding=encoding) as f:
                    text = f.read()
                
                if text.strip():
                    doc = Document(
                        page_content=text,
                        metadata={"source": str(path)}
                    )
                    logger.info(f"Loaded text: {path.name} ({encoding})")
                    return [doc]
                return []
                
            except (UnicodeDecodeError, LookupError):
                continue
            except Exception as e:
                logger.error(f"Failed to load text: {e}")
                return []
        
        logger.error(f"Failed to decode: {path.name}")
        return []
    
    @staticmethod
    def _load_powerpoint(path: Path) -> List[Document]:
        """Load PowerPoint file"""
        try:
            from pptx import Presentation
            
            prs = Presentation(str(path))
            documents = []
            
            for i, slide in enumerate(prs.slides):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                
                if text_parts:
                    text = "\n".join(text_parts)
                    doc = Document(
                        page_content=text,
                        metadata={"source": str(path)}
                    )
                    documents.append(doc)
            
            logger.info(f"Loaded PowerPoint: {path.name} ({len(documents)} slides)")
            return documents
            
        except Exception as e:
            logger.error(f"Failed to load PowerPoint: {e}")
            return []
    
    @staticmethod
    def _load_json(path: Path) -> List[Document]:
        """Load JSON file"""
        try:
            import json
            
            with open(path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            text = json.dumps(data, indent=2, ensure_ascii=False)
            
            doc = Document(
                page_content=text,
                metadata={"source": str(path)}
            )
            
            logger.info(f"Loaded JSON: {path.name}")
            return [doc]
            
        except Exception as e:
            logger.error(f"Failed to load JSON: {e}")
            return []
    
    @staticmethod
    def _load_image(path: Path) -> List[Document]:
        """Load image file (OCR)"""
        try:
            # OCR은 선택적 기능
            logger.warning(f"OCR not implemented for: {path.name}")
            
            doc = Document(
                page_content=f"[Image: {path.name}]",
                metadata={"source": str(path)}
            )
            
            return [doc]
            
        except Exception as e:
            logger.error(f"Failed to load image: {e}")
            return []
    

                        continue
            
            ole.close()
            
            if text_parts:
                full_text = '\n'.join(text_parts)
                doc = Document(
                    page_content=full_text,
                    metadata={"source": str(path)}
                )
                logger.info(f"Loaded HWP: {path.name}")
                return [doc]
            
            logger.warning(f"No text extracted from HWP: {path.name}")
            return []
            
        except ImportError:
            logger.error("olefile not installed. Install with: pip install olefile")
            return []
        except Exception as e:
            logger.error(f"Failed to load HWP: {e}")
            return []
