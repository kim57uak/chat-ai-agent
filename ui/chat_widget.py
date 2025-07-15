from PyQt6.QtWidgets import QWidget, QVBoxLayout, QTextEdit, QLineEdit, QPushButton, QHBoxLayout, QFileDialog, QCheckBox, QLabel, QProgressBar, QTextBrowser, QPlainTextEdit, QComboBox
from PyQt6.QtWebEngineWidgets import QWebEngineView
from PyQt6.QtCore import Qt, pyqtSignal, QObject, QTimer
from PyQt6.QtGui import QFont, QKeySequence, QShortcut
from core.file_utils import load_config, load_model_api_key, load_last_model
from core.ai_client import AIClient
from core.conversation_history import ConversationHistory
from ui.intelligent_formatter import IntelligentContentFormatter
import os
import threading

# 파일 내용 추출 유틸리티
from PyPDF2 import PdfReader
from docx import Document

class AIProcessor(QObject):
    finished = pyqtSignal(str, str, list)  # sender, text, used_tools
    error = pyqtSignal(str)
    streaming = pyqtSignal(str, str)  # sender, partial_text
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self._cancelled = False
    
    def cancel(self):
        self._cancelled = True
    
    def _simulate_streaming(self, sender, response):
        """가상 스트림 출력 시뮬레이션 - 단순화"""
        import time
        import threading
        
        def stream_chunks():
            # 간단한 청크 분할
            chunk_size = 50
            for i in range(0, len(response), chunk_size):
                if self._cancelled:
                    break
                chunk = response[i:i+chunk_size]
                self.streaming.emit(sender, chunk)
                time.sleep(0.1)
        
        if len(response) > 50:
            threading.Thread(target=stream_chunks, daemon=True).start()
    
    def _split_response_into_chunks(self, response):
        """응답을 의미 단위로 분할"""
        import re
        
        # 문장 단위로 분할
        sentences = re.split(r'([.!?]\s+)', response)
        chunks = []
        current_chunk = ""
        
        for i, part in enumerate(sentences):
            current_chunk += part
            
            # 문장 끝이거나 충분히 길면 청크 완성
            if (part.strip().endswith(('.', '!', '?')) or len(current_chunk) > 20) and current_chunk.strip():
                chunks.append(current_chunk)
                current_chunk = ""
        
        # 남은 내용 추가
        if current_chunk.strip():
            chunks.append(current_chunk)
        
        return chunks if chunks else [response]
    
    def process_request(self, api_key, model, messages, user_text=None, agent_mode=False, file_prompt=None):
        """AI 요청 처리 - 메인 스레드에서 실행"""
        
        def _process():
            try:
                if self._cancelled:
                    return
                
                client = AIClient(api_key, model)
                # 대화 히스토리를 클라이언트에 설정
                if messages:
                    client.conversation_history = messages
                    print(f"[디버그] 히스토리 설정: {len(messages)}개")
                
                response = None
                sender = 'AI'
                used_tools = []
                
                if file_prompt:
                    if agent_mode:
                        print(f"[DEBUG] 파일 프롬프트 에이전트 모드로 처리")
                        result = client.agent_chat(file_prompt)
                        if isinstance(result, tuple):
                            response, used_tools = result
                            print(f"[DEBUG] 파일 프롬프트 에이전트 응답 완료, 사용된 도구: {used_tools}")
                        else:
                            response = result
                            used_tools = []
                        sender = '에이전트'
                    else:
                        print(f"[DEBUG] 파일 프롬프트 단순 채팅 모드로 처리")
                        response = client.simple_chat(file_prompt)
                        sender = 'AI'
                        used_tools = []
                else:
                    if agent_mode:
                        print(f"[DEBUG] 에이전트 모드로 처리: {user_text[:50]}...")
                        result = client.agent_chat(user_text)
                        if isinstance(result, tuple):
                            response, used_tools = result
                            print(f"[DEBUG] 에이전트 응답 완료, 사용된 도구: {used_tools}")
                        else:
                            response = result
                            used_tools = []
                        sender = '에이전트'
                    else:
                        print(f"[DEBUG] 단순 채팅 모드로 처리: {user_text[:50]}...")
                        response = client.simple_chat(user_text)
                        sender = 'AI'
                        used_tools = []
                
                if not self._cancelled and response:
                    # AI 응답 길이 디버그
                    print(f"[DEBUG] AI 응답 생성 완룈 - 길이: {len(response)}자")
                    # 안전한 문자열 출력 (공백 제거 및 줄바꿈 변환)
                    safe_start = response[:200].replace('\n', ' ').replace('\r', ' ').strip()
                    print(f"[DEBUG] 응답 내용 시작: {safe_start}...")
                    if len(response) > 500:
                        # 끝부분에서 실제 내용이 있는 부분 찾기
                        trimmed_response = response.rstrip()
                        safe_end = trimmed_response[-200:].replace('\n', ' ').replace('\r', ' ').strip()
                        print(f"[DEBUG] 응답 내용 끝: ...{safe_end}")
                    
                    # 스트리밍 없이 즉시 완성된 응답 표시
                    self.finished.emit(sender, response, used_tools)
                elif not self._cancelled:
                    self.error.emit("응답을 생성할 수 없습니다.")
                    
            except Exception as e:
                if not self._cancelled:
                    self.error.emit(f'오류 발생: {str(e)}')
        
        # 별도 스레드에서 실행
        thread = threading.Thread(target=_process, daemon=True)
        thread.start()

class ChatWidget(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setStyleSheet("background-color: #1a1a1a;")
        self.layout = QVBoxLayout(self)
        self.layout.setContentsMargins(10, 10, 10, 10)
        
        # 대화 히스토리 관리
        self.conversation_history = ConversationHistory()
        self.conversation_history.load_from_file()
        
        # 업로드된 파일 정보 저장
        self.uploaded_file_content = None
        self.uploaded_file_name = None

        # 상단 정보 영역 (모델명 + 도구 상태)
        info_layout = QHBoxLayout()
        
        # 현재 모델명 표시 라벨 (클릭 가능)
        self.model_label = QLabel(self)
        self.model_label.setAlignment(Qt.AlignmentFlag.AlignLeft)
        self.model_label.setCursor(Qt.CursorShape.PointingHandCursor)
        self.model_label.setStyleSheet("""
            QLabel {
                color: rgb(163,135,215);
                font-size: 12px;
                font-weight: bold;
                padding: 10px;
                background-color: #1a1a1a;
            }
            QLabel:hover {
                background-color: #2a2a2a;
                border-radius: 4px;
            }
        """)
        self.model_label.mousePressEvent = self.show_model_popup
        
        # 도구 상태 표시 라벨 (클릭 가능)
        self.tools_label = QLabel(self)
        self.tools_label.setAlignment(Qt.AlignmentFlag.AlignRight)
        self.tools_label.setCursor(Qt.CursorShape.PointingHandCursor)
        self.tools_label.setStyleSheet("""
            QLabel {
                color: rgb(135,163,215);
                font-size: 12px;
                font-weight: bold;
                padding: 10px;
                background-color: #1a1a1a;
            }
            QLabel:hover {
                background-color: #2a2a2a;
                border-radius: 4px;
            }
        """)
        self.tools_label.mousePressEvent = self.show_tools_popup
        
        info_layout.addWidget(self.model_label, 1)
        info_layout.addWidget(self.tools_label, 0)
        self.layout.addLayout(info_layout)
        
        self.update_model_label()
        # 도구 라벨 초기 설정 - 항상 보이도록
        self.tools_label.setText('🔧 도구 확인중...')
        self.tools_label.setVisible(True)
        self.update_tools_label()
        
        # 도구 상태 주기적 갱신 타이머 (초기 지연 후 시작)
        self.tools_update_timer = QTimer()
        self.tools_update_timer.timeout.connect(self.update_tools_label)
        # 초기 업데이트 후 주기적 갱신 시작
        QTimer.singleShot(2000, self.update_tools_label)  # 2초 후 첫 업데이트
        QTimer.singleShot(5000, lambda: self.tools_update_timer.start(10000))  # 5초 후 시작, 10초마다 갱신

        # 채팅 표시 영역 - QWebEngineView로 교체
        self.chat_display = QWebEngineView(self)
        self.chat_display.setMinimumHeight(400)
        self.chat_messages = []  # HTML 메시지 저장용
        
        # QWebEngineView 초기 HTML 설정
        self.init_web_view()
        
        self.layout.addWidget(self.chat_display, 1)
        
        # 로딩 표시
        self.loading_bar = QProgressBar(self)
        self.loading_bar.setRange(0, 0)
        self.loading_bar.hide()
        self.loading_bar.setStyleSheet("""
            QProgressBar {
                background-color: #2a2a2a;
                border: 1px solid #444444;
                border-radius: 4px;
                text-align: center;
                color: #ffffff;
            }
            QProgressBar::chunk {
                background-color: rgb(163,135,215);
                border-radius: 3px;
            }
        """)
        self.layout.addWidget(self.loading_bar)

        input_layout = QHBoxLayout()
        
        # 입력창 컸테이너 (모드 선택과 함께)
        input_container = QWidget(self)
        input_container_layout = QHBoxLayout(input_container)
        input_container_layout.setContentsMargins(0, 0, 0, 0)
        input_container_layout.setSpacing(0)
        
        # 모드 선택 콤보박스 (입력창 내부 왼쪽)
        self.mode_combo = QComboBox(self)
        self.mode_combo.addItems(["Ask", "Agent"])
        self.mode_combo.setCurrentText("Ask")
        self.mode_combo.setStyleSheet("""
            QComboBox {
                background-color: transparent;
                color: #888888;
                border: none;
                border-right: 1px solid #444444;
                border-radius: 0px;
                padding: 4px 8px;
                font-size: 11px;
                font-weight: 500;
                min-width: 45px;
                max-width: 45px;
                outline: none;
            }
            QComboBox:focus {
                border: none;
                border-right: 1px solid #444444;
                outline: none;
            }
            QComboBox:hover {
                color: #ffffff;
                background-color: rgba(255,255,255,0.05);
            }
            QComboBox::drop-down {
                border: none;
                width: 12px;
            }
            QComboBox::down-arrow {
                image: none;
                border-left: 3px solid transparent;
                border-right: 3px solid transparent;
                border-top: 3px solid #888888;
                margin-right: 3px;
            }
            QComboBox:hover::down-arrow {
                border-top: 3px solid #ffffff;
            }
            QComboBox QAbstractItemView {
                background-color: #2a2a2a;
                color: #ffffff;
                selection-background-color: rgb(163,135,215);
                border: 1px solid #444444;
                border-radius: 4px;
            }
        """)
        
        # 입력창
        self.input_text = QTextEdit(self)
        self.input_text.setMaximumHeight(80)
        self.input_text.setPlaceholderText("메시지를 입력하세요... (Ctrl+Enter로 전송)")
        self.input_text.setStyleSheet("""
            QTextEdit {
                background-color: transparent;
                color: #ffffff;
                border: none;
                font-size: 12px;
                padding: 10px;
            }
        """)
        
        # 컸테이너 스타일
        input_container.setStyleSheet("""
            QWidget {
                background-color: #2a2a2a;
                border: 1px solid #444444;
                border-radius: 6px;
            }
        """)
        
        # 컸테이너에 위젯 추가
        input_container_layout.addWidget(self.mode_combo, 0)
        input_container_layout.addWidget(self.input_text, 1)
        
        # 모드 변경 시 플레이스홀더 업데이트
        self.mode_combo.currentTextChanged.connect(self.update_placeholder)
        
        self.send_button = QPushButton('전송', self)
        self.send_button.setMinimumHeight(80)
        self.send_button.setStyleSheet("""
            QPushButton {
                background-color: rgb(163,135,215);
                color: #ffffff;
                border: none;
                border-radius: 6px;
                font-weight: bold;
                font-size: 14px;
            }
            QPushButton:hover {
                background-color: rgb(143,115,195);
            }
            QPushButton:pressed {
                background-color: rgb(123,95,175);
            }
            QPushButton:disabled {
                background-color: #555555;
                color: #888888;
            }
        """)
        
        self.cancel_button = QPushButton('취소', self)
        self.cancel_button.setMinimumHeight(80)
        self.cancel_button.setVisible(False)
        self.cancel_button.setStyleSheet("""
            QPushButton {
                background-color: #F44336;
                color: #ffffff;
                border: none;
                border-radius: 6px;
                font-weight: bold;
                font-size: 12px;
            }
            QPushButton:hover {
                background-color: #D32F2F;
            }
            QPushButton:pressed {
                background-color: #B71C1C;
            }
        """)
        
        self.upload_button = QPushButton('파일\n업로드', self)
        self.upload_button.setMinimumHeight(80)
        self.upload_button.setStyleSheet("""
            QPushButton {
                background-color: rgb(135,163,215);
                color: #ffffff;
                border: none;
                border-radius: 6px;
                font-weight: bold;
                font-size: 14px;
            }
            QPushButton:hover {
                background-color: rgb(115,143,195);
            }
            QPushButton:pressed {
                background-color: rgb(95,123,175);
            }
            QPushButton:disabled {
                background-color: #555555;
                color: #888888;
            }
        """)
        
        input_layout.addWidget(input_container, 5)
        input_layout.addWidget(self.send_button, 1)
        input_layout.addWidget(self.cancel_button, 1)
        input_layout.addWidget(self.upload_button, 1)

        self.layout.addLayout(input_layout, 0)

        self.send_button.clicked.connect(self.send_message)
        self.cancel_button.clicked.connect(self.cancel_request)
        self.upload_button.clicked.connect(self.upload_file)
        
        # Enter키로 전송 (Shift+Enter는 줄바꿈)
        self.input_text.keyPressEvent = self.handle_input_key_press
        
        # Ctrl+Enter로도 전송 가능
        send_shortcut = QShortcut(QKeySequence("Ctrl+Return"), self.input_text)
        send_shortcut.activated.connect(self.send_message)

        self.messages = []
        self.ai_processor = AIProcessor(self)
        self.request_start_time = None
        
        # AI 프로세서 시그널 연결
        self.ai_processor.finished.connect(self.on_ai_response)
        self.ai_processor.error.connect(self.on_ai_error)
        self.ai_processor.streaming.connect(self.on_ai_streaming)
        
        # 타이핑 애니메이션용
        self.typing_timer = QTimer()
        self.typing_timer.timeout.connect(self.show_next_line)
        self.typing_lines = []
        self.current_line_index = 0
        self.current_sender = ""
        self.current_message_id = ""
        self.is_typing = False
        
        # 웹뷰 로드 완료 후 이전 대화 로드
        self.chat_display.loadFinished.connect(self._on_webview_loaded)
        
        # 초기화 완료 후 바로 이전 대화 로드 시도
        QTimer.singleShot(1000, self._load_previous_conversations)
    
    def _safe_run_js(self, js_code):
        """JavaScript를 안전하게 실행"""
        try:
            if len(js_code) > 50000:
                print(f"JavaScript 코드가 너무 김: {len(js_code)}자")
                return
            
            self.chat_display.page().runJavaScript(js_code)
        except Exception as e:
            print(f"JavaScript 실행 오류: {e}")
    
    def handle_input_key_press(self, event):
        """입력창 키 이벤트 처리"""
        from PyQt6.QtCore import Qt
        from PyQt6.QtGui import QKeyEvent
        
        if event.key() == Qt.Key.Key_Return or event.key() == Qt.Key.Key_Enter:
            if event.modifiers() == Qt.KeyboardModifier.ShiftModifier:
                # Shift+Enter: 줄바꿈 삽입
                QTextEdit.keyPressEvent(self.input_text, event)
            else:
                # Enter: 메시지 전송
                self.send_message()
        else:
            # 다른 키들은 기본 처리
            QTextEdit.keyPressEvent(self.input_text, event)
    
    def update_placeholder(self):
        """모드에 따라 플레이스홀더 업데이트"""
        current_mode = self.mode_combo.currentText()
        if current_mode == "Ask":
            self.input_text.setPlaceholderText("단순 질의를 입력하세요... (Enter로 전송, Shift+Enter로 줄바꿈)")
        else:
            self.input_text.setPlaceholderText("도구 사용 가능한 메시지를 입력하세요... (Enter로 전송, Shift+Enter로 줄바꿈)")

    def init_web_view(self):
        """웹 브라우저 초기화"""
        html_template = """
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <style>
                * {
                    box-sizing: border-box;
                }
                
                body {
                    background-color: #1a1a1a;
                    color: #e8e8e8;
                    font-family: -apple-system, BlinkMacSystemFont, 'SF Pro Display', 'Segoe UI', 'Roboto', 'Helvetica Neue', Arial, sans-serif;
                    font-size: 14px;
                    line-height: 1.6;
                    margin: 16px;
                    padding: 0;
                    word-wrap: break-word;
                    overflow-wrap: break-word;
                    overflow-y: auto;
                    height: auto;
                    min-height: 100vh;
                }
                
                /* 코드 블록 스타일 */
                pre {
                    background: #1e1e1e;
                    color: #f8f8f2;
                    padding: 20px;
                    border-radius: 8px;
                    font-family: 'JetBrains Mono', 'Fira Code', 'SF Mono', Consolas, 'Liberation Mono', Menlo, Monaco, monospace;
                    font-size: 13px;
                    line-height: 1.5;
                    overflow-x: auto;
                    white-space: pre;
                    tab-size: 4;
                    border: 1px solid #444;
                    box-shadow: 0 2px 8px rgba(0,0,0,0.3);
                }
                
                /* 인라인 코드 */
                code {
                    background-color: #2d2d2d;
                    color: #f8f8f2;
                    padding: 4px 8px;
                    border-radius: 4px;
                    font-family: 'JetBrains Mono', 'Fira Code', 'SF Mono', Consolas, monospace;
                    font-size: 12px;
                    border: 1px solid #444;
                }
                
                /* 헤딩 스타일 */
                h1, h2, h3, h4, h5, h6 {
                    margin-top: 24px;
                    margin-bottom: 12px;
                    font-weight: 600;
                    line-height: 1.25;
                }
                
                h1 { font-size: 24px; color: #ffffff; border-bottom: 2px solid #444; padding-bottom: 8px; }
                h2 { font-size: 20px; color: #eeeeee; border-bottom: 1px solid #333; padding-bottom: 6px; }
                h3 { font-size: 18px; color: #dddddd; }
                h4 { font-size: 16px; color: #cccccc; }
                h5 { font-size: 14px; color: #bbbbbb; }
                h6 { font-size: 13px; color: #aaaaaa; }
                
                /* 링크 스타일 */
                a {
                    color: #87CEEB;
                    text-decoration: none;
                    border-bottom: 1px dotted #87CEEB;
                    transition: all 0.2s ease;
                }
                
                a:hover {
                    color: #B0E0E6;
                    border-bottom: 1px solid #B0E0E6;
                }
                
                /* 리스트 스타일 */
                ul, ol {
                    padding-left: 20px;
                    margin: 12px 0;
                }
                
                li {
                    margin: 4px 0;
                    color: #cccccc;
                }
                
                /* 인용문 */
                blockquote {
                    margin: 16px 0;
                    padding: 12px 16px;
                    border-left: 4px solid #87CEEB;
                    background-color: rgba(135, 206, 235, 0.1);
                    color: #dddddd;
                    font-style: italic;
                }
                
                /* 테이블 스타일 */
                table {
                    border-collapse: collapse;
                    width: auto;
                    margin: 16px 0;
                    background-color: #2a2a2a;
                    border-radius: 8px;
                    overflow: hidden;
                    box-shadow: 0 4px 12px rgba(0,0,0,0.3);
                }
                
                th, td {
                    padding: 12px 16px;
                    text-align: left;
                    border: 1px solid #444;
                    white-space: normal;
                    word-wrap: break-word;
                    vertical-align: top;
                }
                
                th {
                    background: linear-gradient(135deg, #3a3a3a, #4a4a4a);
                    color: #ffffff;
                    font-weight: 700;
                    font-size: 13px;
                    text-transform: uppercase;
                    letter-spacing: 0.5px;
                }
                
                tr:nth-child(even) {
                    background-color: #252525;
                }
                
                tr:hover {
                    background-color: #333333;
                }
                
                /* 수평선 */
                hr {
                    border: none;
                    height: 2px;
                    background: linear-gradient(to right, transparent, #444, transparent);
                    margin: 20px 0;
                }
                
                /* 강조 텍스트 */
                strong {
                    color: #ffffff;
                    font-weight: 600;
                }
                
                em {
                    color: #dddddd;
                    font-style: italic;
                }
                
                del {
                    color: #888888;
                    text-decoration: line-through;
                }
                
                /* 메시지 컸테이너 */
                .message {
                    margin: 16px 0;
                    padding: 16px;
                    border-radius: 12px;
                    box-shadow: 0 4px 12px rgba(0,0,0,0.3);
                }
                
                .user { background: rgba(163,135,215,0.15); border-left: 4px solid rgb(163,135,215); }
                .ai { background: rgba(135,163,215,0.15); border-left: 4px solid rgb(135,163,215); }
                .system { background: rgba(215,163,135,0.15); border-left: 4px solid rgb(215,163,135); }
                
                /* 스크롤바 스타일링 */
                ::-webkit-scrollbar {
                    width: 8px;
                    height: 8px;
                }
                
                ::-webkit-scrollbar-track {
                    background: #2a2a2a;
                    border-radius: 4px;
                }
                
                ::-webkit-scrollbar-thumb {
                    background: #555;
                    border-radius: 4px;
                }
                
                ::-webkit-scrollbar-thumb:hover {
                    background: #666;
                }
            </style>
            <script>
                function copyCode(codeId) {
                    const codeElement = document.getElementById(codeId);
                    const text = codeElement.textContent;
                    
                    if (navigator.clipboard && window.isSecureContext) {
                        navigator.clipboard.writeText(text).then(() => {
                            showCopySuccess(event.target);
                        }).catch(err => {
                            fallbackCopy(text, event.target);
                        });
                    } else {
                        fallbackCopy(text, event.target);
                    }
                }
                
                function fallbackCopy(text, button) {
                    const textarea = document.createElement('textarea');
                    textarea.value = text;
                    textarea.style.position = 'fixed';
                    textarea.style.opacity = '0';
                    document.body.appendChild(textarea);
                    textarea.select();
                    
                    try {
                        document.execCommand('copy');
                        showCopySuccess(button);
                    } catch (err) {
                        button.textContent = '실패';
                        button.style.background = '#F44336';
                        setTimeout(() => {
                            button.textContent = '복사';
                            button.style.background = '#444';
                        }, 1500);
                    } finally {
                        document.body.removeChild(textarea);
                    }
                }
                
                function showCopySuccess(button) {
                    button.textContent = '복사됨!';
                    button.style.background = '#4CAF50';
                    setTimeout(() => {
                        button.textContent = '복사';
                        button.style.background = '#444';
                    }, 1500);
                }
            </script>
        </head>
        <body>
            <div id="messages"></div>
        </body>
        </html>
        """
        self.chat_display.setHtml(html_template)
    
    def update_model_label(self):
        model = load_last_model()
        self.model_label.setText(f'현재 모델: <b>{model}</b> 📋')
    
    def show_model_popup(self, event):
        """사용 가능한 모델 목록 팝업 표시"""
        try:
            from PyQt6.QtWidgets import QMenu
            from core.file_utils import load_config, save_last_model
            
            config = load_config()
            models = config.get('models', {})
            
            if not models:
                return
            
            menu = QMenu(self)
            menu.setStyleSheet("""
                QMenu {
                    background-color: #2a2a2a;
                    color: #ffffff;
                    border: 1px solid #444444;
                    border-radius: 4px;
                    padding: 4px;
                }
                QMenu::item {
                    padding: 8px 16px;
                    border-radius: 2px;
                }
                QMenu::item:selected {
                    background-color: rgb(163,135,215);
                }
            """)
            
            current_model = load_last_model()
            
            for model_name, model_config in models.items():
                if model_config.get('api_key'):  # API 키가 있는 모델만 표시
                    action = menu.addAction(f"🤖 {model_name}")
                    if model_name == current_model:
                        action.setText(f"✅ {model_name} (현재)")
                    action.triggered.connect(lambda checked, m=model_name: self.change_model(m))
            
            # 라벨 위치에서 팝업 표시
            menu.exec(self.model_label.mapToGlobal(event.pos()))
            
        except Exception as e:
            print(f"모델 팝업 표시 오류: {e}")
    
    def change_model(self, model_name):
        """모델 변경"""
        try:
            from core.file_utils import save_last_model
            save_last_model(model_name)
            self.update_model_label()
            self.append_chat('시스템', f'모델이 {model_name}으로 변경되었습니다.')
            print(f"[디버그] 모델 변경: {model_name}")
        except Exception as e:
            print(f"모델 변경 오류: {e}")
            self.append_chat('시스템', f'모델 변경 중 오류가 발생했습니다: {e}')
    
    def update_tools_label(self):
        """활성화된 도구 수 표시 업데이트 - 동기 처리"""
        try:
            from mcp.servers.mcp import get_all_mcp_tools
            tools = get_all_mcp_tools()
            tool_count = len(tools) if tools else 0
            
            if tool_count > 0:
                text = f'🔧 {tool_count}개 도구 활성화'
            else:
                text = '🔧 도구 없음'
            
            self.tools_label.setText(text)
            print(f"[디버그] 도구 라벨 업데이트: {text}")
            
        except Exception as e:
            self.tools_label.setText('🔧 도구 상태 불명')
            print(f"도구 라벨 업데이트 오류: {e}")
    

    

    
    def show_tools_popup(self, event):
        """활성화된 도구 목록 팝업 표시 - 개선된 처리"""
        print("[디버그] 도구 팝업 클릭됨")
        
        try:
            # 직접 동기 방식으로 처리 (간단하게)
            from PyQt6.QtWidgets import QMenu, QMessageBox
            from mcp.servers.mcp import get_all_mcp_tools
            
            print("[디버그] 도구 목록 조회 시작")
            tools = get_all_mcp_tools()
            print(f"[디버그] 조회된 도구 수: {len(tools) if tools else 0}")
            
            if not tools:
                # 도구가 없을 때 메시지 표시
                QMessageBox.information(
                    self, 
                    "도구 상태", 
                    "활성화된 MCP 도구가 없습니다.\n\n설정 > MCP 서버 관리에서 서버를 활성화하세요."
                )
                return
            
            # 메뉴 생성 및 표시
            self._show_tools_menu(event, tools)
            
        except Exception as e:
            print(f"[디버그] 도구 팝업 표시 오류: {e}")
            # 폴백: 간단한 메시지 표시
            try:
                from PyQt6.QtWidgets import QMessageBox
                QMessageBox.warning(
                    self, 
                    "오류", 
                    f"도구 상태를 확인할 수 없습니다.\n\n오류: {e}"
                )
            except:
                print("[디버그] 폴백 메시지 표시도 실패")
    
    def _show_tools_menu(self, event, tools):
        """도구 메뉴 표시 (메인 스레드에서 실행)"""
        try:
            print(f"[디버그] 메뉴 생성 시작, 도구 수: {len(tools)}")
            from PyQt6.QtWidgets import QMenu
            
            menu = QMenu(self)
            menu.setStyleSheet("""
                QMenu {
                    background-color: #2a2a2a;
                    color: #ffffff;
                    border: 1px solid #444444;
                    border-radius: 4px;
                    padding: 4px;
                }
                QMenu::item {
                    padding: 6px 12px;
                    border-radius: 2px;
                }
                QMenu::item:selected {
                    background-color: #444444;
                }
            """)
            
            # 서버별로 도구 그룹화
            servers = {}
            for tool in tools:
                if isinstance(tool, str):
                    tool_name = tool
                    server_name = 'Tools'
                else:
                    # MCPTool 객체의 속성 접근
                    tool_name = tool.name if hasattr(tool, 'name') else str(tool)
                    server_name = tool.server_name if hasattr(tool, 'server_name') else 'Tools'
                
                if server_name not in servers:
                    servers[server_name] = []
                servers[server_name].append(tool_name)
            
            print(f"[디버그] 서버별 그룹화 완료: {list(servers.keys())}")
            
            # 메뉴 항목 추가
            for server_name, tool_names in servers.items():
                menu.addAction(f"📦 {server_name} ({len(tool_names)}개)")
                for tool_name in tool_names[:5]:  # 최대 5개만 표시
                    menu.addAction(f"  • {tool_name}")
                if len(tool_names) > 5:
                    menu.addAction(f"  ... 외 {len(tool_names)-5}개")
                menu.addSeparator()
            
            print("[디버그] 메뉴 표시 시작")
            # 마우스 커서 위치에 표시
            from PyQt6.QtGui import QCursor
            menu.exec(QCursor.pos())
            print("[디버그] 메뉴 표시 완료")
            
        except Exception as e:
            print(f"[디버그] 도구 메뉴 표시 오류: {e}")
            import traceback
            traceback.print_exc()

    def send_message(self):
        user_text = self.input_text.toPlainText().strip()
        if not user_text:
            return
        
        # 이전 요청이 진행 중이면 취소
        if hasattr(self, 'ai_processor'):
            self.ai_processor.cancel()
            self.ai_processor = AIProcessor(self)
            self.ai_processor.finished.connect(self.on_ai_response)
            self.ai_processor.error.connect(self.on_ai_error)
            self.ai_processor.streaming.connect(self.on_ai_streaming)
        
        self._process_new_message(user_text)
    
    def _process_new_message(self, user_text):
        """새 메시지 처리"""
        from datetime import datetime
        self.request_start_time = datetime.now()
        
        self.append_chat('사용자', user_text)
        self.input_text.clear()
        
        # 히스토리에 사용자 메시지 추가
        self.conversation_history.add_message('user', user_text)
        self.conversation_history.save_to_file()
        self.messages.append({'role': 'user', 'content': user_text})

        model = load_last_model()
        api_key = load_model_api_key(model)
        self.update_model_label()
        if not api_key:
            self.append_chat('시스템', 'API Key가 설정되어 있지 않습니다. 환경설정에서 입력해 주세요.')
            return
        
        # 업로드된 파일이 있으면 파일 내용과 함께 분석
        if self.uploaded_file_content:
            # 이미지 파일인지 확인
            if "[IMAGE_BASE64]" in self.uploaded_file_content:
                # 이미지의 경우 사용자 요청과 이미지 데이터를 직접 결합
                combined_prompt = f'{user_text}\n\n{self.uploaded_file_content}'
                print(f"[디버그] 이미지 프롬프트 생성, 길이: {len(combined_prompt)}")
                print(f"[디버그] 종료 태그 확인: {'[/IMAGE_BASE64]' in combined_prompt}")
            else:
                # 일반 파일의 경우 기존 방식 유지
                combined_prompt = f'업로드된 파일 ({self.uploaded_file_name})에 대한 사용자 요청: {user_text}\n\n파일 내용:\n{self.uploaded_file_content}'
            
            self._start_ai_request(api_key, model, None, combined_prompt)
            # 파일 내용 초기화
            self.uploaded_file_content = None
            self.uploaded_file_name = None
        else:
            self._start_ai_request(api_key, model, user_text)
    
    def _start_ai_request(self, api_key, model, user_text, file_prompt=None):
        """모드에 따라 AI 요청 시작"""
        self.set_ui_enabled(False)
        self.show_loading(True)
        
        # 최신 대화 히스토리 가져오기 (파일에서 다시 로드)
        self.conversation_history.load_from_file()
        recent_history = self.conversation_history.get_recent_messages(10)
        print(f"[디버그] 전달할 히스토리: {len(recent_history)}개")
        for i, msg in enumerate(recent_history[-3:]):
            content = msg.get('content', '')[:50].replace('\n', ' ').replace('\r', ' ').strip()
            print(f"  [{i}] {msg.get('role', 'unknown')}: {content}...")
        
        # 모드에 따라 에이전트 사용 여부 결정
        current_mode = self.mode_combo.currentText()
        use_agent = (current_mode == "Agent")
        print(f"[DEBUG] 선택된 모드: {current_mode}, 에이전트 사용: {use_agent}")
        
        self.ai_processor.process_request(
            api_key, model, recent_history, user_text,
            use_agent, file_prompt
        )
    
    def old_start_ai_request(self, api_key, model, user_text, file_prompt=None):
        """AI 요청 시작 - 단순화된 방식"""
        self.set_ui_enabled(False)
        self.show_loading(True)
        
        # 최근 대화 기록 포함하여 전송
        recent_history = self.conversation_history.get_recent_messages(10)
        
        # AI 요청 처리 - 항상 에이전트 모드
        self.ai_processor.process_request(
            api_key, model, recent_history, user_text,
            True, file_prompt
        )
        
        # 타임아웃 제거 - 사용자가 취소 버튼으로 제어
    


    def upload_file(self):
        file_path, _ = QFileDialog.getOpenFileName(self, '파일 선택', '', '모든 파일 (*);;텍스트 파일 (*.txt);;PDF 파일 (*.pdf);;Word 파일 (*.docx *.doc);;Excel 파일 (*.xlsx *.xls);;PowerPoint 파일 (*.pptx *.ppt);;JSON 파일 (*.json);;이미지 파일 (*.jpg *.jpeg *.png *.gif *.bmp *.webp);;CSV 파일 (*.csv)')
        if not file_path:
            return
        
        # 이전 요청이 진행 중이면 취소
        if hasattr(self, 'ai_processor'):
            self.ai_processor.cancel()
            self.ai_processor = AIProcessor(self)
            self.ai_processor.finished.connect(self.on_ai_response)
            self.ai_processor.error.connect(self.on_ai_error)
            self.ai_processor.streaming.connect(self.on_ai_streaming)
        
        self._process_file_upload(file_path)
    
    def _process_file_upload(self, file_path):
        """파일 업로드 처리"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            content = ""
            
            if ext == '.txt':
                # 다양한 encoding 시도
                for encoding in ['utf-8', 'cp949', 'euc-kr', 'latin-1', 'utf-8-sig']:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            content = f.read()
                        break
                    except UnicodeDecodeError:
                        continue
                else:
                    content = f"텍스트 파일: {os.path.basename(file_path)}\n인코딩을 읽을 수 없습니다."
            elif ext == '.pdf':
                reader = PdfReader(file_path)
                content = "\n".join(page.extract_text() or '' for page in reader.pages)
            elif ext in ['.docx', '.doc']:
                doc = Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs])
            elif ext in ['.xlsx', '.xls']:
                try:
                    import pandas as pd
                    df = pd.read_excel(file_path)
                    content = f"파일 정보: {os.path.basename(file_path)}\n열 수: {len(df.columns)}, 행 수: {len(df)}\n\n데이터 미리보기:\n{df.head(10).to_string()}"
                except ImportError:
                    content = f"Excel 파일: {os.path.basename(file_path)}\n(pandas 라이브러리가 필요합니다. pip install pandas openpyxl)"
            elif ext in ['.pptx', '.ppt']:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    slides_text = []
                    for i, slide in enumerate(prs.slides, 1):
                        slide_text = f"슬라이드 {i}:\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text += shape.text + "\n"
                        slides_text.append(slide_text)
                    content = "\n\n".join(slides_text)
                except ImportError:
                    content = f"PowerPoint 파일: {os.path.basename(file_path)}\n(python-pptx 라이브러리가 필요합니다. pip install python-pptx)"
            elif ext == '.json':
                import json
                # JSON 파일도 다중 인코딩 지원
                for encoding in ['utf-8', 'utf-8-sig', 'cp949', 'euc-kr', 'latin-1']:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            data = json.load(f)
                        content = f"JSON 파일 내용:\n{json.dumps(data, indent=2, ensure_ascii=False)}"
                        break
                    except (UnicodeDecodeError, json.JSONDecodeError):
                        continue
                else:
                    content = f"JSON 파일: {os.path.basename(file_path)}\n파일을 읽을 수 없습니다."
            elif ext == '.csv':
                try:
                    import pandas as pd
                    # 다양한 encoding 시도
                    for encoding in ['utf-8', 'cp949', 'euc-kr', 'latin-1']:
                        try:
                            df = pd.read_csv(file_path, encoding=encoding)
                            content = f"CSV 파일 정보: {os.path.basename(file_path)}\n열 수: {len(df.columns)}, 행 수: {len(df)}\n\n데이터 미리보기:\n{df.head(10).to_string()}"
                            break
                        except UnicodeDecodeError:
                            continue
                except ImportError:
                    # pandas 없을 때 기본 처리
                    for encoding in ['utf-8', 'cp949', 'euc-kr', 'latin-1']:
                        try:
                            with open(file_path, 'r', encoding=encoding) as f:
                                lines = f.readlines()[:10]
                                content = f"CSV 파일: {os.path.basename(file_path)}\n미리보기 (10줄):\n{''.join(lines)}"
                                break
                        except UnicodeDecodeError:
                            continue
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']:
                import base64
                try:
                    # 이미지를 base64로 인코딩
                    with open(file_path, 'rb') as img_file:
                        img_data = base64.b64encode(img_file.read()).decode('utf-8')
                    
                    # 이미지 정보 추가 (선택사항)
                    img_info = ""
                    try:
                        from PIL import Image
                        img = Image.open(file_path)
                        img_info = f"\n이미지 정보: {img.size[0]}x{img.size[1]} 픽셀, 모드: {img.mode}"
                    except (ImportError, ModuleNotFoundError):
                        pass
                    
                    content = f"[IMAGE_BASE64]{img_data}[/IMAGE_BASE64]\n이미지 파일: {os.path.basename(file_path)}{img_info}"
                    print(f"[디버그] 이미지 콘텐츠 생성 완료, 길이: {len(content)}")
                    print(f"[디버그] 종료 태그 위치: {content.rfind('[/IMAGE_BASE64]')}")
                    
                except Exception as e:
                    file_size = os.path.getsize(file_path)
                    content = f"이미지 파일: {os.path.basename(file_path)}\n파일 크기: {file_size:,} bytes\n이미지 처리 오류: {str(e)}"
            else:
                # 지원하지 않는 파일도 기본 텍스트로 읽기 시도
                for encoding in ['utf-8', 'cp949', 'euc-kr', 'latin-1']:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            content = f.read()
                        break
                    except UnicodeDecodeError:
                        continue
                    except Exception as e:
                        content = f"파일: {os.path.basename(file_path)}\n파일 읽기 오류: {str(e)}"
                        break
            
            self.append_chat('사용자', f'📎 파일 업로드: {os.path.basename(file_path)}')
            
            # 이미지 파일의 경우 자르지 않음 (태그 보존)
            if "[IMAGE_BASE64]" not in content and len(content) > 5000:
                content = content[:5000] + "...(내용 생략)"
            
            # 파일 내용을 임시 저장
            self.uploaded_file_content = content
            self.uploaded_file_name = os.path.basename(file_path)
            print(f"[디버그] 업로드된 파일 내용 길이: {len(content)}")
            if "[IMAGE_BASE64]" in content:
                print(f"[디버그] 이미지 데이터 포함 확인")
                print(f"[디버그] 종료 태그 위치: {content.rfind('[/IMAGE_BASE64]')}")
            
            # 사용자에게 프롬프트 입력 안내
            self.append_chat('시스템', f'파일이 업로드되었습니다. 이제 파일에 대해 무엇을 알고 싶은지 메시지를 입력해주세요.')
            self.input_text.setPlaceholderText(f"{self.uploaded_file_name}에 대해 무엇을 알고 싶으신가요? (Enter로 전송)")
            
        except Exception as e:
            self.append_chat('시스템', f'파일 처리 오류: {e}')
            # 오류 시 파일 내용 초기화
            self.uploaded_file_content = None
            self.uploaded_file_name = None
            self.input_text.setPlaceholderText("메시지를 입력하세요... (Enter로 전송, Shift+Enter로 줄바꿈)")

    def cancel_request(self):
        """요청 취소 - 단순화된 방식"""
        print("취소 요청 시작")
        
        # UI 상태 복원
        self.set_ui_enabled(True)
        self.show_loading(False)
        
        # AI 프로세서 취소
        if hasattr(self, 'ai_processor'):
            self.ai_processor.cancel()
        
        self.append_chat('시스템', '요청을 취소했습니다.')
        print("취소 요청 완료")
    
    def on_ai_streaming(self, sender, partial_text):
        """스트림 출력 처리"""
        if not hasattr(self, 'current_stream_message_id'):
            import uuid
            self.current_stream_message_id = f"stream_{uuid.uuid4().hex[:8]}"
            self.current_stream_content = ""
            self._create_stream_message_container(sender, self.current_stream_message_id)
        
        self.current_stream_content += partial_text
        
        # 기본 HTML 이스케이프
        formatted_content = self.current_stream_content.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br>')
        formatted_content = formatted_content.replace('\\', '\\\\').replace("'", "\\'")
        
        js_code = f"""
        try {{
            var contentDiv = document.getElementById('{self.current_stream_message_id}_content');
            if (contentDiv) {{
                contentDiv.innerHTML = '{formatted_content}';
                window.scrollTo(0, document.body.scrollHeight);
            }}
        }} catch(e) {{}}
        """
        self._safe_run_js(js_code)
    
    def _create_stream_message_container(self, sender, message_id):
        """스트림 메시지 컸테이너 생성"""
        # 발신자별 스타일
        if sender == '사용자':
            bg_color = 'rgba(163,135,215,0.15)'
            border_color = 'rgb(163,135,215)'
            icon = '💬'
            sender_color = 'rgb(163,135,215)'
        elif sender in ['AI', '에이전트'] or '에이전트' in sender:
            bg_color = 'rgba(135,163,215,0.15)'
            border_color = 'rgb(135,163,215)'
            icon = '🤖'
            sender_color = 'rgb(135,163,215)'
        else:
            bg_color = 'rgba(215,163,135,0.15)'
            border_color = 'rgb(215,163,135)'
            icon = '⚙️'
            sender_color = 'rgb(215,163,135)'
        
        html_container = f"""
        <div id="{message_id}" style="
            margin: 12px 0;
            padding: 16px;
            background: linear-gradient(135deg, {bg_color}33, {bg_color}11);
            border-radius: 12px;
            border-left: 4px solid {border_color};
            box-shadow: 0 4px 12px rgba(0,0,0,0.3);
        ">
            <div style="
                margin: 0 0 12px 0;
                font-weight: 700;
                color: {sender_color};
                font-size: 12px;
                display: flex;
                align-items: center;
                gap: 8px;
                text-transform: uppercase;
                letter-spacing: 0.5px;
            ">
                <span style="font-size: 16px;">{icon}</span>
                <span>{sender}</span>
                <span style="opacity: 0.6; font-size: 10px;">• 입력 중...</span>
            </div>
            <div id="{message_id}_content" style="
                margin: 0;
                padding-left: 24px;
                line-height: 1.6;
                color: #ffffff;
                font-size: 13px;
                word-wrap: break-word;
                font-family: 'SF Pro Display', 'Segoe UI', Arial, sans-serif;
            ">
                <span style="opacity: 0.5;">●</span>
            </div>
        </div>
        """
        
        # JavaScript로 메시지 추가
        html_escaped = html_container.replace('\\', '\\\\').replace('`', '\\`').replace('${', '\\${')
        js_code = f"""
        try {{
            var messagesDiv = document.getElementById('messages');
            var messageDiv = document.createElement('div');
            messageDiv.innerHTML = `{html_escaped}`;
            messagesDiv.appendChild(messageDiv);
            window.scrollTo(0, document.body.scrollHeight);
        }} catch(e) {{
            console.log('Message add error:', e);
        }}
        """
        self._safe_run_js(js_code)
    
    def on_ai_response(self, sender, text, used_tools):
        # AI 응답 길이 디버그 (안전한 출력)
        print(f"[DEBUG] AI 응답 받음 - 길이: {len(text)}자")
        safe_start = text[:100].replace('\n', ' ').replace('\r', ' ').strip()
        # 끝부분에서 실제 내용이 있는 부분 찾기
        trimmed_text = text.rstrip()
        safe_end = trimmed_text[-100:].replace('\n', ' ').replace('\r', ' ').strip()
        print(f"[DEBUG] 응답 시작: {safe_start}...")
        print(f"[DEBUG] 응답 끝: ...{safe_end}")
        
        # 테이블 후처리 적용
        processed_text = self._post_process_tables(text)
        
        # 응답 시간 계산
        response_time = ""
        if self.request_start_time:
            from datetime import datetime
            elapsed = datetime.now() - self.request_start_time
            response_time = f" ({elapsed.total_seconds():.1f}초)"
        
        # 현재 모델 정보 가져오기
        current_model = load_last_model()
        
        # 모델명과 응답시간을 응답 끝에 추가
        enhanced_text = f"{processed_text}\n\n---\n*🤖 {current_model}{response_time}*"
        
        # 스트리밍 없이 즉시 완성된 응답 표시
        self.append_chat(sender, enhanced_text)
        
        # 히스토리에는 원본 텍스트만 저장
        self.conversation_history.add_message('assistant', text)
        self.conversation_history.save_to_file()
        
        self.messages.append({'role': 'assistant', 'content': text})
        self.set_ui_enabled(True)
        self.show_loading(False)
    
    def _get_tool_emoji(self, used_tools):
        """사용된 도구에 따라 이모티콘 반환 (동적 매핑)"""
        if not used_tools:
            return ""
        
        # 도구 이름 키워드 기반 이모티콘 매핑
        emoji_map = {
            'search': '🔍',
            'web': '🌐', 
            'url': '🌐',
            'fetch': '📄',
            'database': '🗄️',
            'mysql': '🗄️',
            'sql': '🗄️',
            'travel': '✈️',
            'tour': '✈️',
            'hotel': '🏨',
            'flight': '✈️',
            'map': '🗺️',
            'location': '📍',
            'geocode': '📍',
            'weather': '🌤️',
            'email': '📧',
            'file': '📁',
            'excel': '📊',
            'chart': '📈',
            'image': '🖼️',
            'translate': '🌐',
            'api': '🔧'
        }
        
        # 첫 번째 도구의 이름에서 키워드 찾기
        tool_name = str(used_tools[0]).lower() if used_tools else ""
        
        for keyword, emoji in emoji_map.items():
            if keyword in tool_name:
                return emoji
        
        # 매핑되지 않은 도구의 기본 이모티콘
        return "⚡"

    def on_ai_error(self, msg):
        # 오류 시에도 응답 시간 표시
        error_time = ""
        if self.request_start_time:
            from datetime import datetime
            elapsed = datetime.now() - self.request_start_time
            error_time = f" (오류발생시간: {elapsed.total_seconds():.1f}초)"
        
        self.append_chat('시스템', msg + error_time)
        self.set_ui_enabled(True)
        self.show_loading(False)

    def set_ui_enabled(self, enabled):
        self.send_button.setEnabled(enabled)
        self.cancel_button.setVisible(not enabled)
        self.upload_button.setEnabled(enabled)
    
    def show_loading(self, show):
        """로딩 상태 표시/숨김"""
        if show:
            self.loading_bar.show()
        else:
            self.loading_bar.hide()

    def _append_simple_chat(self, sender, text):
        """간단한 채팅 메시지 표시"""
        self.append_chat(sender, text)
    
    def append_chat(self, sender, text):
        """채팅 메시지 표시 - 안정화"""
        # 발신자별 스타일
        if sender == '사용자':
            bg_color = 'rgba(163,135,215,0.15)'
            border_color = 'rgb(163,135,215)'
            icon = '💬'
            sender_color = 'rgb(163,135,215)'
        elif sender in ['AI', '에이전트'] or '에이전트' in sender:
            bg_color = 'rgba(135,163,215,0.15)'
            border_color = 'rgb(135,163,215)'
            icon = '🤖'
            sender_color = 'rgb(135,163,215)'
        else:
            bg_color = 'rgba(215,163,135,0.15)'
            border_color = 'rgb(215,163,135)'
            icon = '⚙️'
            sender_color = 'rgb(215,163,135)'
        
        # 마크다운 처리 - IntelligentContentFormatter 사용
        formatter = IntelligentContentFormatter()
        formatted_text = formatter.format_content(text)
        
        # Base64 인코딩으로 안전하게 전달
        import base64
        import uuid
        message_id = f"msg_{uuid.uuid4().hex[:8]}"
        
        # 1단계: 메시지 컨테이너 생성
        create_js = f'''
        try {{
            var messagesDiv = document.getElementById('messages');
            var messageDiv = document.createElement('div');
            messageDiv.id = '{message_id}';
            messageDiv.style.cssText = 'margin:12px 0;padding:16px;background:{bg_color};border-radius:12px;border-left:4px solid {border_color};';
            
            var headerDiv = document.createElement('div');
            headerDiv.style.cssText = 'margin:0 0 12px 0;font-weight:700;color:{sender_color};font-size:12px;display:flex;align-items:center;gap:8px;';
            headerDiv.innerHTML = '<span style="font-size:16px;">{icon}</span><span>{sender}</span>';
            
            var contentDiv = document.createElement('div');
            contentDiv.id = '{message_id}_content';
            contentDiv.style.cssText = 'margin:0;padding-left:24px;line-height:1.6;color:#ffffff;font-size:13px;word-wrap:break-word;';
            
            messageDiv.appendChild(headerDiv);
            messageDiv.appendChild(contentDiv);
            messagesDiv.appendChild(messageDiv);
            
            window.scrollTo(0, document.body.scrollHeight);
        }} catch(e) {{
            console.log('Create message error:', e);
        }}
        '''
        
        # 2단계: 콘텐츠 설정 - JSON 안전 전달 방식
        def set_content():
            import json
            # JSON.stringify를 사용하여 모든 특수문자를 안전하게 전달
            safe_content = json.dumps(formatted_text, ensure_ascii=False)
            
            content_js = f'''
            try {{
                var contentDiv = document.getElementById('{message_id}_content');
                if (contentDiv) {{
                    contentDiv.innerHTML = {safe_content};
                    window.scrollTo(0, document.body.scrollHeight);
                }}
            }} catch(e) {{
                console.log('Set content error:', e);
                var contentDiv = document.getElementById('{message_id}_content');
                if (contentDiv) {{
                    contentDiv.textContent = 'Content display error';
                }}
            }}
            '''
            self.chat_display.page().runJavaScript(content_js)
        
        self.chat_display.page().runJavaScript(create_js)
        # 짧은 지연 후 콘텐츠 설정
        QTimer.singleShot(50, set_content)
    
    def start_optimized_typing(self, sender, text):
        """즉시 메시지 표시"""
        self.append_chat(sender, text)
    
    def show_next_line(self):
        """다음 줄 표시"""
        if self.current_line_index >= len(self.typing_lines):
            self.typing_timer.stop()
            self.is_typing = False
            return
            
        line = self.typing_lines[self.current_line_index]
        formatted_line = self.format_text(line)
        
        # 현재 메시지 컨테이너에 줄 추가
        line_escaped = formatted_line.replace('\\', '\\\\').replace('`', '\\`').replace('${', '\\${')
        line_js = f"""
        try {{
            var contentDiv = document.getElementById('{self.current_message_id}_content');
            var lineDiv = document.createElement('div');
            lineDiv.innerHTML = `{line_escaped}` + '<br>';
            contentDiv.appendChild(lineDiv);
            window.scrollTo(0, document.body.scrollHeight);
        }} catch(e) {{
            console.log('Line add error:', e);
        }}
        """
        
        self._safe_run_js(line_js)
        self.current_line_index += 1
    
    def _split_text_for_typing(self, text):
        """타이핑용 텍스트 분할"""
        lines = text.split('\n')
        chunks = []
        current_chunk = []
        current_length = 0
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if len(line) > 100 or '```' in line or line.startswith('#'):
                if current_chunk:
                    chunks.append('\n'.join(current_chunk))
                    current_chunk = []
                    current_length = 0
                chunks.append(line)
            else:
                current_chunk.append(line)
                current_length += len(line)
                
                if current_length > 200:
                    chunks.append('\n'.join(current_chunk))
                    current_chunk = []
                    current_length = 0
        
        if current_chunk:
            chunks.append('\n'.join(current_chunk))
            
        return chunks
    
    def format_text(self, text):
        """Simple text formatting without LLM"""
        return self._basic_format_text(text)
    

    

    

    

    
    def _post_process_tables(self, text):
        """테이블 후처리 - 구분선 정규화 및 긴 셀 처리"""
        import re
        
        if '|' not in text or '---' not in text:
            return text
        
        lines = text.split('\n')
        processed_lines = []
        
        for line in lines:
            # 테이블 구분선 감지 및 정규화
            if '|' in line and ('---' in line or ':--' in line or '--:' in line):
                # 구분선 정규화
                parts = line.split('|')
                normalized_parts = []
                
                for part in parts:
                    part = part.strip()
                    if part and ('-' in part or ':' in part):
                        normalized_parts.append('---')
                    else:
                        normalized_parts.append(part)
                
                processed_lines.append('|'.join(normalized_parts))
            elif '|' in line and line.count('|') >= 2:
                # 일반 테이블 행 - 긴 셀 처리
                parts = line.split('|')
                wrapped_parts = []
                
                for part in parts:
                    part = part.strip()
                    if len(part) > 30:
                        # 긴 텍스트 줄바꿈
                        wrapped = self._wrap_long_text(part, 30)
                        wrapped_parts.append(wrapped)
                    else:
                        wrapped_parts.append(part)
                
                processed_lines.append('|'.join(wrapped_parts))
            else:
                processed_lines.append(line)
        
        return '\n'.join(processed_lines)
    
    def _wrap_long_text(self, text, max_width):
        """긴 텍스트 자동 줄바꿈 - 비정상적으로 긴 텍스트 처리"""
        if len(text) <= max_width:
            return text
        
        # 비정상적으로 긴 텍스트 (오류 데이터) 처리
        if len(text) > 1000:
            return f'<span style="color:#ff6b6b;font-style:italic;">[{len(text)}자 데이터 - 표시 생략]</span>'
        
        # 공백 기준 분할 시도
        words = text.split(' ')
        if len(words) > 1:
            lines = []
            current_line = []
            current_length = 0
            
            for word in words:
                if current_length + len(word) + 1 <= max_width:
                    current_line.append(word)
                    current_length += len(word) + 1
                else:
                    if current_line:
                        lines.append(' '.join(current_line))
                    current_line = [word[:max_width]]  # 너무 긴 단어 자르기
                    current_length = len(current_line[0])
            
            if current_line:
                lines.append(' '.join(current_line))
            
            # 최대 5줄로 제한
            if len(lines) > 5:
                lines = lines[:5]
                lines.append('<span style="color:#888;">... (내용 생략)</span>')
            
            return '<br>'.join(lines)
        
        # 강제 분할 (최대 200자로 제한)
        if len(text) > 200:
            return text[:200] + '<br><span style="color:#888;">... (내용 생략)</span>'
        
        chunks = [text[i:i+max_width] for i in range(0, len(text), max_width)]
        return '<br>'.join(chunks)
    

    
    def scroll_to_bottom(self):
        """스크롤을 맨 아래로"""
        self.chat_display.verticalScrollBar().setValue(
            self.chat_display.verticalScrollBar().maximum()
        )

    def _on_webview_loaded(self, ok):
        """웹뷰 로드 완료 후 이전 대화 로드"""
        if ok:
            QTimer.singleShot(500, self._load_previous_conversations)
    
    def _load_previous_conversations(self):
        """이전 대화 내용 로드 - 원본 그대로"""
        try:
            self.conversation_history.load_from_file()
            recent_messages = self.conversation_history.get_recent_messages(3)
            
            if recent_messages:
                self._append_simple_chat('시스템', f'이전 대화 {len(recent_messages)}개를 불러왔습니다.')
                
                for msg in recent_messages:
                    role = msg.get('role', '')
                    content = msg.get('content', '')
                    
                    if not content or not content.strip():
                        continue
                    
                    # 내용 생략하지 않음
                    
                    if role == 'user':
                        self._append_simple_chat('사용자', content)
                    elif role == 'assistant':
                        self._append_simple_chat('AI', content)
            else:
                self._append_simple_chat('시스템', '새로운 대화를 시작합니다.')
                
        except Exception as e:
            self._append_simple_chat('시스템', '새로운 대화를 시작합니다.')
    
    def clear_conversation_history(self):
        """대화 히스토리 초기화"""
        self.conversation_history.clear_session()
        self.conversation_history.save_to_file()
        self.messages = []
        print("[디버그] 대화 히스토리가 초기화되었습니다.")
        
        # 웹뷰도 초기화
        self.init_web_view()
    
    def close(self):
        """위젯 종료 - 단순화된 방식"""
        print("ChatWidget 종료 시작")
        
        try:
            # 타이머 정지
            if hasattr(self, 'typing_timer') and self.typing_timer:
                self.typing_timer.stop()
            
            # AI 프로세서 취소
            if hasattr(self, 'ai_processor'):
                self.ai_processor.cancel()
            
            print("ChatWidget 종료 완료")
            
        except Exception as e:
            print(f"ChatWidget 종료 중 오류: {e}")
        
        # 타이머 정지
        if hasattr(self, 'tools_update_timer'):
            self.tools_update_timer.stop()