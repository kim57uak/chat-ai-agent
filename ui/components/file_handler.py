import os
import base64
from PyPDF2 import PdfReader
from docx import Document


class FileHandler:
    """파일 처리를 담당하는 클래스 (SRP)"""
    
    @staticmethod
    def process_file(file_path):
        """파일 내용 추출"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            content = ""
            
            if ext == '.txt':
                content = FileHandler._read_text_file(file_path)
            elif ext == '.pdf':
                content = FileHandler._read_pdf_file(file_path)
            elif ext in ['.docx', '.doc']:
                content = FileHandler._read_word_file(file_path)
            elif ext in ['.xlsx', '.xls']:
                content = FileHandler._read_excel_file(file_path)
            elif ext in ['.pptx', '.ppt']:
                content = FileHandler._read_powerpoint_file(file_path)
            elif ext == '.json':
                content = FileHandler._read_json_file(file_path)
            elif ext == '.csv':
                content = FileHandler._read_csv_file(file_path)
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']:
                content = FileHandler._read_image_file(file_path)
            else:
                content = FileHandler._read_text_file(file_path)
            
            return content, os.path.basename(file_path)
            
        except Exception as e:
            return f"파일 처리 오류: {str(e)}", os.path.basename(file_path)
    
    @staticmethod
    def _read_text_file(file_path):
        """텍스트 파일 읽기"""
        for encoding in ['utf-8', 'cp949', 'euc-kr', 'latin-1', 'utf-8-sig']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        return f"텍스트 파일: {os.path.basename(file_path)}\n인코딩을 읽을 수 없습니다."
    
    @staticmethod
    def _read_pdf_file(file_path):
        """PDF 파일 읽기"""
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or '' for page in reader.pages)
    
    @staticmethod
    def _read_word_file(file_path):
        """Word 파일 읽기"""
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    
    @staticmethod
    def _read_excel_file(file_path):
        """Excel 파일 읽기"""
        try:
            import pandas as pd
            df = pd.read_excel(file_path)
            return f"파일 정보: {os.path.basename(file_path)}\n열 수: {len(df.columns)}, 행 수: {len(df)}\n\n데이터 미리보기:\n{df.head(10).to_string()}"
        except ImportError:
            return f"Excel 파일: {os.path.basename(file_path)}\n(pandas 라이브러리가 필요합니다. pip install pandas openpyxl)"
    
    @staticmethod
    def _read_powerpoint_file(file_path):
        """PowerPoint 파일 읽기"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"슬라이드 {i}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                slides_text.append(slide_text)
            return "\n\n".join(slides_text)
        except ImportError:
            return f"PowerPoint 파일: {os.path.basename(file_path)}\n(python-pptx 라이브러리가 필요합니다. pip install python-pptx)"
    
    @staticmethod
    def _read_json_file(file_path):
        """JSON 파일 읽기"""
        import json
        for encoding in ['utf-8', 'utf-8-sig', 'cp949', 'euc-kr', 'latin-1']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    data = json.load(f)
                return f"JSON 파일 내용:\n{json.dumps(data, indent=2, ensure_ascii=False)}"
            except (UnicodeDecodeError, json.JSONDecodeError):
                continue
        return f"JSON 파일: {os.path.basename(file_path)}\n파일을 읽을 수 없습니다."
    
    @staticmethod
    def _read_csv_file(file_path):
        """CSV 파일 읽기"""
        try:
            import pandas as pd
            for encoding in ['utf-8', 'cp949', 'euc-kr', 'latin-1']:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    return f"CSV 파일 정보: {os.path.basename(file_path)}\n열 수: {len(df.columns)}, 행 수: {len(df)}\n\n데이터 미리보기:\n{df.head(10).to_string()}"
                except UnicodeDecodeError:
                    continue
        except ImportError:
            for encoding in ['utf-8', 'cp949', 'euc-kr', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        lines = f.readlines()[:10]
                        return f"CSV 파일: {os.path.basename(file_path)}\n미리보기 (10줄):\n{''.join(lines)}"
                except UnicodeDecodeError:
                    continue
        return f"CSV 파일을 읽을 수 없습니다: {os.path.basename(file_path)}"
    
    @staticmethod
    def _read_image_file(file_path):
        """이미지 파일 읽기"""
        try:
            with open(file_path, 'rb') as img_file:
                img_data = base64.b64encode(img_file.read()).decode('utf-8')
            
            img_info = ""
            try:
                from PIL import Image
                img = Image.open(file_path)
                img_info = f"\n이미지 정보: {img.size[0]}x{img.size[1]} 픽셀, 모드: {img.mode}"
            except (ImportError, ModuleNotFoundError):
                pass
            
            return f"[IMAGE_BASE64]{img_data}[/IMAGE_BASE64]\n이미지 파일: {os.path.basename(file_path)}{img_info}"
            
        except Exception as e:
            file_size = os.path.getsize(file_path)
            return f"이미지 파일: {os.path.basename(file_path)}\n파일 크기: {file_size:,} bytes\n이미지 처리 오류: {str(e)}"